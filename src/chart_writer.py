"""Native PowerPoint pie chart with theme-coloured slices and readable legend.

We patch the chart XML directly because python-pptx's font wrappers write
`<a:defRPr>` only — and PowerPoint's chart engine often overrides defaults
from the embedded chart style. Writing the values directly into the
`<c:txPr>` elements (and forcing `<c:legendPos val="r"/>`) ensures the
sizes, weights, and colours actually render.
"""
from __future__ import annotations

from typing import Sequence

from lxml import etree
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

from src import theme


_C_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def build_pie_chart(
    slide,
    x,
    y,
    cx,
    cy,
    labels: Sequence[str],
    values: Sequence[float],
    slice_colors: Sequence[RGBColor],
) -> None:
    chart_data = CategoryChartData()
    chart_data.categories = list(labels)
    chart_data.add_series("", list(values))

    graphic = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data)
    chart = graphic.chart

    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    _set_legend_manual_layout(
        chart,
        x=theme.LEGEND_LAYOUT_X,
        y=theme.LEGEND_LAYOUT_Y,
        w=theme.LEGEND_LAYOUT_W,
        h=theme.LEGEND_LAYOUT_H,
    )

    # Per-slice fills.
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.show_percentage = True
    dl.show_category_name = False
    dl.show_value = False
    dl.position = XL_LABEL_POSITION.OUTSIDE_END
    dl.number_format = "0.0%"
    dl.number_format_is_linked = False

    series = plot.series[0]
    for idx, point in enumerate(series.points):
        fill = point.format.fill
        fill.solid()
        fill.fore_color.rgb = slice_colors[idx % len(slice_colors)]

    # --- Direct XML patches (PowerPoint ignores defRPr from chart style) ---
    chart_space = chart._chartSpace

    # 1. Force <c:legendPos val="r"/> — python-pptx writes it without a val
    legend_pos = chart_space.find(
        f"{{{_C_NS}}}chart/{{{_C_NS}}}legend/{{{_C_NS}}}legendPos"
    )
    if legend_pos is not None:
        legend_pos.set("val", "r")

    # 2. Force the chart style to a neutral one so the chart's theme XML
    #    doesn't override our text settings. Insert <c:style val="2"/>.
    _ensure_chart_style(chart_space, val=2)

    # 3. Rewrite legend txPr with explicit rPr (not defRPr) — runs win over
    #    inheritance. Apply size/bold/color/font + Sinhala cs.
    legend = chart_space.find(f"{{{_C_NS}}}chart/{{{_C_NS}}}legend")
    if legend is not None:
        _force_text_style(
            legend,
            size_pt=theme.CHART_LEGEND_SIZE.pt,
            bold=True,
            color_hex=_hex(theme.BLACK),
            sinhala=True,
        )

    # 4. Same for data label txPr — bold percentages, larger size.
    pie = chart_space.find(
        f"{{{_C_NS}}}chart/{{{_C_NS}}}plotArea/{{{_C_NS}}}pieChart"
    )
    if pie is not None:
        for d_lbls in pie.iter(f"{{{_C_NS}}}dLbls"):
            _force_text_style(
                d_lbls,
                size_pt=theme.CHART_DATA_LABEL_SIZE.pt,
                bold=True,
                color_hex=_hex(theme.BLACK),
                sinhala=False,
            )


# --- XML helpers -----------------------------------------------------------

def _set_legend_manual_layout(chart, x: float, y: float, w: float, h: float) -> None:
    """Force a wide legend container via `<c:manualLayout>` (fractions 0..1).

    python-pptx's auto-positioned legend collapses to ~20% width — too narrow
    for long Sinhala category labels at 36pt. Source slide 3 used 43% width,
    so we set an explicit layout that mirrors that.
    """
    legend = chart._chartSpace.find(
        f"{{{_C_NS}}}chart/{{{_C_NS}}}legend"
    )
    if legend is None:
        return
    existing = legend.find(f"{{{_C_NS}}}layout")
    if existing is not None:
        legend.remove(existing)
    layout = etree.Element(f"{{{_C_NS}}}layout")
    manual = etree.SubElement(layout, f"{{{_C_NS}}}manualLayout")
    etree.SubElement(manual, f"{{{_C_NS}}}xMode", val="edge")
    etree.SubElement(manual, f"{{{_C_NS}}}yMode", val="edge")
    etree.SubElement(manual, f"{{{_C_NS}}}x", val=str(x))
    etree.SubElement(manual, f"{{{_C_NS}}}y", val=str(y))
    etree.SubElement(manual, f"{{{_C_NS}}}w", val=str(w))
    etree.SubElement(manual, f"{{{_C_NS}}}h", val=str(h))
    # `<c:layout>` must come BEFORE `<c:overlay>` per the schema
    legend_pos = legend.find(f"{{{_C_NS}}}legendPos")
    if legend_pos is not None:
        legend_pos.addnext(layout)
    else:
        legend.insert(0, layout)


def _ensure_chart_style(chart_space, val: int) -> None:
    chart = chart_space.find(f"{{{_C_NS}}}chart")
    if chart is None:
        return
    style = chart_space.find(f"{{{_C_NS}}}style")
    if style is None:
        style = etree.SubElement(chart_space, f"{{{_C_NS}}}style")
        # <c:style> must come AFTER <c:chart> per ECMA-376; ensure ordering
        chart_space.remove(style)
        chart.addnext(style)
    style.set("val", str(val))


def _force_text_style(host_element, size_pt: float, bold: bool,
                      color_hex: str, sinhala: bool) -> None:
    """Replace `<c:txPr>` under `host_element` with one that uses an explicit
    rPr (forces PowerPoint to honour our settings)."""
    txPr_qn = f"{{{_C_NS}}}txPr"
    existing = host_element.find(txPr_qn)
    if existing is not None:
        host_element.remove(existing)

    txPr = etree.SubElement(host_element, txPr_qn)
    etree.SubElement(txPr, f"{{{_A_NS}}}bodyPr")
    etree.SubElement(txPr, f"{{{_A_NS}}}lstStyle")
    p = etree.SubElement(txPr, f"{{{_A_NS}}}p")
    pPr = etree.SubElement(p, f"{{{_A_NS}}}pPr")

    sz_hundredths = int(size_pt * 100)
    defRPr = etree.SubElement(pPr, f"{{{_A_NS}}}defRPr",
                              sz=str(sz_hundredths),
                              b="1" if bold else "0")
    solid = etree.SubElement(defRPr, f"{{{_A_NS}}}solidFill")
    etree.SubElement(solid, f"{{{_A_NS}}}srgbClr", val=color_hex)
    etree.SubElement(defRPr, f"{{{_A_NS}}}latin",
                     typeface=theme.SINHALA_FONT if sinhala else theme.LATIN_FONT)
    etree.SubElement(defRPr, f"{{{_A_NS}}}cs",
                     typeface=theme.SINHALA_FONT)

    # Add an empty run with explicit rPr — matches what PowerPoint expects.
    endRPr = etree.SubElement(p, f"{{{_A_NS}}}endParaRPr",
                              lang="en-US",
                              sz=str(sz_hundredths),
                              b="1" if bold else "0")
    solid2 = etree.SubElement(endRPr, f"{{{_A_NS}}}solidFill")
    etree.SubElement(solid2, f"{{{_A_NS}}}srgbClr", val=color_hex)


def _hex(color: RGBColor) -> str:
    return str(color)


# ===========================================================================
# v2 chart helpers — bars, lines, stacked bars
# ===========================================================================


def build_bar_chart(
    slide,
    x, y, cx, cy,
    categories: Sequence[str],
    series_data: dict[str, Sequence[float]],
    palette: Sequence[RGBColor],
) -> None:
    """Clustered (grouped) column chart. `series_data = {label: [values...]}`."""
    _build_categorical(
        slide, x, y, cx, cy,
        chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED,
        categories=categories,
        series_data=series_data,
        palette=palette,
        show_data_labels=True,
    )


def build_line_chart(
    slide,
    x, y, cx, cy,
    categories: Sequence[str],
    series_data: dict[str, Sequence[float]],
    palette: Sequence[RGBColor],
) -> None:
    """Line chart, single or multi-series."""
    _build_categorical(
        slide, x, y, cx, cy,
        chart_type=XL_CHART_TYPE.LINE_MARKERS,
        categories=categories,
        series_data=series_data,
        palette=palette,
        show_data_labels=True,
        line_chart=True,
    )


def build_stacked_bar_chart(
    slide,
    x, y, cx, cy,
    categories: Sequence[str],
    series_data: dict[str, Sequence[float]],
    palette: Sequence[RGBColor],
) -> None:
    """Stacked column chart — one stack per category, segments per series."""
    _build_categorical(
        slide, x, y, cx, cy,
        chart_type=XL_CHART_TYPE.COLUMN_STACKED,
        categories=categories,
        series_data=series_data,
        palette=palette,
        show_data_labels=False,  # too crowded inside segments
    )


def _build_categorical(
    slide, x, y, cx, cy,
    chart_type,
    categories: Sequence[str],
    series_data: dict[str, Sequence[float]],
    palette: Sequence[RGBColor],
    show_data_labels: bool,
    line_chart: bool = False,
) -> None:
    chart_data = CategoryChartData()
    chart_data.categories = list(categories)
    for name, values in series_data.items():
        chart_data.add_series(name, list(values))

    graphic = slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data)
    chart = graphic.chart

    chart.has_title = False
    chart.has_legend = len(series_data) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.TOP
        chart.legend.include_in_layout = False

    # Per-series colors
    plot = chart.plots[0]
    for s_idx, series in enumerate(plot.series):
        color = palette[s_idx % len(palette)]
        if line_chart:
            series.format.line.color.rgb = color
            try:
                series.format.line.width = theme.Pt(3)
            except Exception:
                pass
        else:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = color

    # Data labels (where useful)
    if show_data_labels:
        for series in plot.series:
            series.data_labels.show_value = True
            series.data_labels.number_format = '#,##0'
            series.data_labels.number_format_is_linked = False

    # XML patches — same recipe as build_pie_chart
    chart_space = chart._chartSpace
    _ensure_chart_style(chart_space, val=2)
    if chart.has_legend:
        legend = chart_space.find(f"{{{_C_NS}}}chart/{{{_C_NS}}}legend")
        if legend is not None:
            legend_pos = legend.find(f"{{{_C_NS}}}legendPos")
            if legend_pos is not None:
                legend_pos.set("val", "t")
            _force_text_style(
                legend,
                size_pt=theme.CHART_LEGEND_SIZE.pt,
                bold=True,
                color_hex=_hex(theme.BLACK),
                sinhala=True,
            )

    # Axes — bump tick label font; without this matplotlib-style charts come
    # out tiny on a 13" projection.
    _force_axis_text_style(chart_space, theme.CHART_LEGEND_SIZE.pt)

    # Data labels styling
    plot_area = chart_space.find(f"{{{_C_NS}}}chart/{{{_C_NS}}}plotArea")
    if plot_area is not None:
        for d_lbls in plot_area.iter(f"{{{_C_NS}}}dLbls"):
            _force_text_style(
                d_lbls,
                size_pt=theme.CHART_DATA_LABEL_SIZE.pt * 0.6,  # ~22pt for bars/lines
                bold=True,
                color_hex=_hex(theme.BLACK),
                sinhala=False,
            )


def _force_axis_text_style(chart_space, size_pt: float) -> None:
    """Bump the font size on all <c:catAx>/<c:valAx> tick labels."""
    plot_area = chart_space.find(f"{{{_C_NS}}}chart/{{{_C_NS}}}plotArea")
    if plot_area is None:
        return
    for ax in list(plot_area.iter(f"{{{_C_NS}}}catAx")) + list(
        plot_area.iter(f"{{{_C_NS}}}valAx")
    ):
        _force_text_style(
            ax,
            size_pt=size_pt,
            bold=False,
            color_hex=_hex(theme.BLACK),
            sinhala=True,
        )
