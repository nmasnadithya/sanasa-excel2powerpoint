"""Builder — opens the source deck, uses its layouts/master, builds new
slides with positions and fonts calibrated to slides 1-3 of the source.

v2 adds trend visuals: bar/line/stacked-bar charts, YTD/delta/top-N tables,
and matplotlib-rendered images (heatmaps, waterfall, sankey, etc.).
"""
from __future__ import annotations

import shutil
from datetime import date
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches

from src import theme
from src.excel_reader import (
    ExcelReader, Row,
    INCOME_SHEET, EXPENSE_SHEET, SUMMARY_SHEET,
)
from src.sinhala_font import apply_sinhala_font
from src.slide_specs import (
    SlideSpec, DataQuery,
    LOAN_SURPLUS_TITLE_NEGATIVE, LOAN_SURPLUS_TITLE_POSITIVE,
    INCOME_TOTAL_LABEL, EXPENSE_TOTAL_LABEL,
)
from src.chart_writer import (
    build_pie_chart, build_bar_chart, build_line_chart, build_stacked_bar_chart,
)
from src.runtime_paths import app_dir, default_template_path


MAX_ROWS_PER_TABLE_SLIDE = 9
CONTINUATION_SUFFIX = " (අඛණ්ඩව)"


def _default_source() -> Path:
    return default_template_path()


# Backwards-compatible name kept for callers; resolved lazily.
DEFAULT_SOURCE = _default_source()
TMP_DIR = app_dir() / "output" / "_tmp"

_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


class TemplateBuilder:
    """Open the source deck, build new slides on top, then strip originals."""

    def __init__(self, source_path: Path | None = None):
        self.source_path = Path(source_path) if source_path else DEFAULT_SOURCE

    # --- Public API --------------------------------------------------------

    def build(
        self,
        specs: list[SlideSpec],
        reader: ExcelReader,
        target: date,
        output_path: Path,
    ) -> None:
        prs = Presentation(str(self.source_path))
        cover_layout = prs.slides[0].slide_layout
        blank_layout = self._find_layout(prs, "Blank")

        original_sld_ids = list(prs.slides._sldIdLst)
        TMP_DIR.mkdir(parents=True, exist_ok=True)

        try:
            for spec in specs:
                self._dispatch(prs, cover_layout, blank_layout, spec, reader, target)
            self._remove_slides(prs, original_sld_ids)
            output_path.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(output_path))
        finally:
            # Clean up transient PNGs whether or not the save succeeded.
            shutil.rmtree(TMP_DIR, ignore_errors=True)

    def _dispatch(self, prs, cover_layout, blank_layout, spec, reader, target):
        layout = spec.layout
        if layout == "cover":
            self._render_cover(prs, cover_layout, spec)
        elif layout == "table":
            self._render_table(prs, blank_layout, spec, reader, target)
        elif layout == "big_number":
            self._render_big_number(prs, blank_layout, spec, reader, target)
        elif layout == "chart":
            self._render_chart(prs, blank_layout, spec, reader, target)
        elif layout == "bar_compare":
            self._render_bar_compare(prs, blank_layout, spec, reader, target)
        elif layout == "line_trend":
            self._render_line_trend(prs, blank_layout, spec, reader, target)
        elif layout == "stacked_bar":
            self._render_stacked_bar(prs, blank_layout, spec, reader, target)
        elif layout == "ytd_table":
            self._render_ytd_table(prs, blank_layout, spec, reader, target)
        elif layout == "delta_table":
            self._render_delta_table(prs, blank_layout, spec, reader, target)
        elif layout == "top_n_table":
            self._render_top_n_table(prs, blank_layout, spec, reader, target)
        elif layout == "image":
            self._render_image(prs, blank_layout, spec, reader, target)
        else:
            raise ValueError(f"Unknown layout {spec.layout!r}")

    # --- v1 Renderers ------------------------------------------------------

    def _render_cover(self, prs, layout, spec: SlideSpec) -> None:
        slide = prs.slides.add_slide(layout)
        # Cover keeps its layout-defined background image — don't paint over it.

        title_ph, subtitle_ph = self._cover_placeholders(slide)

        if title_ph is not None:
            title_ph.left = theme.COVER_TITLE_LEFT
            title_ph.top = theme.COVER_TITLE_TOP
            title_ph.width = theme.COVER_TITLE_WIDTH
            title_ph.height = theme.COVER_TITLE_HEIGHT
            self._set_textframe_text(
                title_ph.text_frame, spec.title,
                size=None, color=None, bold=True, align=None,
            )

        if subtitle_ph is not None and spec.subtitle:
            subtitle_ph.left = theme.COVER_SUBTITLE_LEFT
            subtitle_ph.top = theme.COVER_SUBTITLE_TOP
            subtitle_ph.width = theme.COVER_SUBTITLE_WIDTH
            subtitle_ph.height = theme.COVER_SUBTITLE_HEIGHT
            self._set_textframe_text(
                subtitle_ph.text_frame, spec.subtitle,
                size=theme.COVER_SUBTITLE_SIZE, color=theme.RED_BRIGHT,
                bold=True, align=None,
            )

        decor = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            theme.COVER_DECOR_LEFT, theme.COVER_DECOR_TOP,
            theme.COVER_DECOR_WIDTH, theme.COVER_DECOR_HEIGHT,
        )
        decor.fill.solid()
        decor.fill.fore_color.rgb = theme.COVER_DECOR_COLOR
        decor.line.fill.background()

    def _render_table(
        self, prs, layout, spec: SlideSpec, reader: ExcelReader, target: date
    ) -> None:
        rows = self._fetch_rows(spec.data, reader, target)

        # v1 retrofit: when surplus_role is set, inject the surplus row + an
        # adjusted-total row. Pagination still applies in case rows + injection
        # exceeds MAX_ROWS_PER_TABLE_SLIDE.
        if spec.surplus_role is not None:
            rows = self._inject_surplus_and_total(rows, spec.surplus_role,
                                                  reader, target)
        if not rows:
            return

        chunks = _distribute_evenly(rows, MAX_ROWS_PER_TABLE_SLIDE)
        for page_idx, chunk in enumerate(chunks):
            slide = prs.slides.add_slide(layout)
            self._set_white_background(slide)
            title = spec.title if page_idx == 0 else spec.title + CONTINUATION_SUFFIX
            self._draw_table_title(slide, title)
            self._draw_table(slide, chunk)

    def _render_big_number(
        self, prs, layout, spec: SlideSpec, reader: ExcelReader, target: date
    ) -> None:
        if spec.computed_key != "loan_surplus":
            raise ValueError(f"Unknown computed_key {spec.computed_key!r}")
        value = abs(reader.loan_surplus(target))
        if value == 0:
            return

        slide = prs.slides.add_slide(layout)
        self._set_white_background(slide)

        title_box = slide.shapes.add_textbox(
            theme.BIG_NUMBER_TITLE_LEFT, theme.BIG_NUMBER_TITLE_TOP,
            theme.BIG_NUMBER_TITLE_WIDTH, theme.BIG_NUMBER_TITLE_HEIGHT,
        )
        self._set_textframe_text(
            title_box.text_frame, spec.title,
            size=theme.BIG_NUMBER_TITLE_SIZE, color=theme.RED,
            bold=False, align=PP_ALIGN.CENTER,
        )

        num_box = slide.shapes.add_textbox(
            theme.BIG_NUMBER_VALUE_LEFT, theme.BIG_NUMBER_VALUE_TOP,
            theme.BIG_NUMBER_VALUE_WIDTH, theme.BIG_NUMBER_VALUE_HEIGHT,
        )
        self._set_textframe_text(
            num_box.text_frame, f"{value:,.2f}",
            size=theme.BIG_NUMBER_SIZE, color=theme.BIG_NUMBER_VALUE_COLOR,
            bold=True, align=PP_ALIGN.CENTER, apply_sinhala=False,
        )

    def _render_chart(
        self, prs, layout, spec: SlideSpec, reader: ExcelReader, target: date
    ) -> None:
        rows = self._fetch_rows(spec.data, reader, target)

        # v1 retrofit: append a surplus slice if surplus sign matches role
        if spec.surplus_role is not None:
            extra = self._surplus_row_for_role(spec.surplus_role, reader, target)
            if extra is not None:
                rows = list(rows) + [extra]

        if not rows:
            return

        slide = prs.slides.add_slide(layout)
        self._set_white_background(slide)
        self._draw_chart_title(slide, spec.title)
        build_pie_chart(
            slide,
            x=theme.CHART_LEFT, y=theme.CHART_TOP,
            cx=theme.CHART_WIDTH, cy=theme.CHART_HEIGHT,
            labels=[r.label for r in rows],
            values=[r.value for r in rows],
            slice_colors=theme.PIE_PALETTE,
        )

    # --- v2 Trend renderers ------------------------------------------------

    def _render_bar_compare(
        self, prs, layout, spec: SlideSpec, reader: ExcelReader, target: date
    ) -> None:
        dates = reader.populated_dates()
        if not dates:
            return
        income = reader.adjusted_monthly_totals(INCOME_SHEET)
        expense = reader.adjusted_monthly_totals(EXPENSE_SHEET)
        slide = prs.slides.add_slide(layout)
        self._set_white_background(slide)
        self._draw_chart_title(slide, spec.title)
        build_bar_chart(
            slide,
            x=theme.CHART_LEFT, y=theme.CHART_TOP,
            cx=theme.CHART_WIDTH, cy=theme.CHART_HEIGHT,
            categories=[_fmt_month(d) for d in dates],
            series_data={
                "ආදායම": [income.get(d, 0.0) for d in dates],
                "වියදම": [expense.get(d, 0.0) for d in dates],
            },
            palette=[theme.LINE_INCOME_COLOR, theme.LINE_EXPENSE_COLOR],
        )

    def _render_line_trend(
        self, prs, layout, spec: SlideSpec, reader: ExcelReader, target: date
    ) -> None:
        dates = reader.populated_dates()
        if not dates:
            return
        slide = prs.slides.add_slide(layout)
        self._set_white_background(slide)
        self._draw_chart_title(slide, spec.title)

        if spec.computed_key == "net_profit":
            net = reader.net_profit_per_month()
            series = {"ශුද්ධ ලාභය": [net.get(d, 0.0) for d in dates]}
            palette = [theme.LINE_PROFIT_COLOR]
        elif spec.computed_key == "cumulative_io":
            income = reader.adjusted_monthly_totals(INCOME_SHEET)
            expense = reader.adjusted_monthly_totals(EXPENSE_SHEET)
            cum_inc, cum_exp, ai, ae = [], [], 0.0, 0.0
            for d in dates:
                ai += income.get(d, 0.0)
                ae += expense.get(d, 0.0)
                cum_inc.append(ai)
                cum_exp.append(ae)
            series = {"සමුච්චිත ආදායම": cum_inc, "සමුච්චිත වියදම": cum_exp}
            palette = [theme.LINE_INCOME_COLOR, theme.LINE_EXPENSE_COLOR]
        else:
            raise ValueError(f"Unknown line_trend computed_key {spec.computed_key!r}")

        build_line_chart(
            slide,
            x=theme.CHART_LEFT, y=theme.CHART_TOP,
            cx=theme.CHART_WIDTH, cy=theme.CHART_HEIGHT,
            categories=[_fmt_month(d) for d in dates],
            series_data=series,
            palette=palette,
        )

    def _render_stacked_bar(
        self, prs, layout, spec: SlideSpec, reader: ExcelReader, target: date
    ) -> None:
        dates = reader.populated_dates()
        if not dates or spec.data is None:
            return
        multi = reader.rows_multi(spec.data.sheet, spec.data.rows,
                                  spec.data.label_col)
        # Build one series per category
        series: dict[str, list[float]] = {}
        for mr in multi:
            series[mr.label] = [mr.values.get(d, 0.0) for d in dates]

        # Add surplus segment when sign matches the role
        if spec.surplus_role is not None:
            surplus = reader.loan_surplus_per_month()
            sign_threshold = (lambda v: v < 0) if spec.surplus_role == "income" \
                else (lambda v: v > 0)
            surplus_series = [
                abs(surplus.get(d, 0.0)) if sign_threshold(surplus.get(d, 0.0)) else 0.0
                for d in dates
            ]
            if any(v > 0 for v in surplus_series):
                label = (LOAN_SURPLUS_TITLE_NEGATIVE if spec.surplus_role == "income"
                         else LOAN_SURPLUS_TITLE_POSITIVE)
                series[label] = surplus_series

        slide = prs.slides.add_slide(layout)
        self._set_white_background(slide)
        self._draw_chart_title(slide, spec.title)
        build_stacked_bar_chart(
            slide,
            x=theme.CHART_LEFT, y=theme.CHART_TOP,
            cx=theme.CHART_WIDTH, cy=theme.CHART_HEIGHT,
            categories=[_fmt_month(d) for d in dates],
            series_data=series,
            palette=theme.STACKED_PALETTE,
        )

    def _render_ytd_table(
        self, prs, layout, spec: SlideSpec, reader: ExcelReader, target: date
    ) -> None:
        if spec.data is None:
            return
        multi = reader.rows_multi(spec.data.sheet, spec.data.rows,
                                  spec.data.label_col)
        if not multi:
            return
        cur_total = sum(mr.ytd_total for mr in multi) or 1.0
        rows: list[list[str]] = []
        for mr in multi:
            cur = mr.values.get(target, 0.0)
            ytd = mr.ytd_total
            pct = (ytd / cur_total) * 100 if cur_total else 0
            rows.append([
                mr.label,
                f"{cur:,.2f}",
                f"{ytd:,.2f}",
                f"{pct:.1f}%",
            ])
        # Total row
        cur_sum = sum(mr.values.get(target, 0.0) for mr in multi)
        ytd_sum = sum(mr.ytd_total for mr in multi)
        rows.append(["එකතුව", f"{cur_sum:,.2f}", f"{ytd_sum:,.2f}", "100.0%"])

        slide = prs.slides.add_slide(layout)
        self._set_white_background(slide)
        self._draw_table_title(slide, spec.title)
        self._draw_n_col_table(
            slide, rows,
            headers=["වර්ගය", _fmt_month(target), "වසරට", "%"],
            col_widths=[
                Emu(4_000_000),
                Emu(2_400_000),
                Emu(2_400_000),
                Emu(1_300_000),
            ],
            col_aligns=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT],
            col_sinhala=[True, False, False, False],
        )

    def _render_delta_table(
        self, prs, layout, spec: SlideSpec, reader: ExcelReader, target: date
    ) -> None:
        if spec.data is None:
            return
        dates = reader.populated_dates()
        prev_date = self._previous_date(dates, target)
        multi = reader.rows_multi(spec.data.sheet, spec.data.rows,
                                  spec.data.label_col)
        if not multi:
            return

        rows: list[list[str]] = []
        delta_colors: list = []
        for mr in multi:
            cur = mr.values.get(target, 0.0)
            prev = mr.values.get(prev_date, 0.0) if prev_date else None
            if prev is None:
                delta_str = "—"
                pct_str = "—"
                color = None
            else:
                delta = cur - prev
                pct = (delta / prev * 100) if prev else 0
                delta_str = f"{delta:+,.2f}"
                pct_str = f"{pct:+.1f}%"
                # Income: positive delta = good (green); Expense: positive = bad (red)
                if spec.compare_kind == "income":
                    color = theme.DELTA_POSITIVE_COLOR if delta >= 0 else theme.DELTA_NEGATIVE_COLOR
                else:  # expense
                    color = theme.DELTA_POSITIVE_COLOR if delta <= 0 else theme.DELTA_NEGATIVE_COLOR
            rows.append([
                mr.label,
                f"{prev:,.2f}" if prev is not None else "—",
                f"{cur:,.2f}",
                delta_str,
                pct_str,
            ])
            delta_colors.append(color)

        slide = prs.slides.add_slide(layout)
        self._set_white_background(slide)
        self._draw_table_title(slide, spec.title)
        self._draw_n_col_table(
            slide, rows,
            headers=[
                "වර්ගය",
                _fmt_month(prev_date) if prev_date else "—",
                _fmt_month(target),
                "වෙනස",
                "වෙනස %",
            ],
            col_widths=[
                Emu(3_400_000),
                Emu(1_900_000),
                Emu(1_900_000),
                Emu(1_500_000),
                Emu(1_400_000),
            ],
            col_aligns=[PP_ALIGN.LEFT] + [PP_ALIGN.RIGHT] * 4,
            col_sinhala=[True, False, False, False, False],
            colorize_columns={3: delta_colors, 4: delta_colors},
        )

    def _render_top_n_table(
        self, prs, layout, spec: SlideSpec, reader: ExcelReader, target: date
    ) -> None:
        if spec.data is None:
            return
        n = spec.top_n or 5
        items = reader.top_line_items(spec.data.sheet, target, n=n)
        if not items:
            return
        slide = prs.slides.add_slide(layout)
        self._set_white_background(slide)
        self._draw_table_title(slide, spec.title)
        self._draw_table(slide, items)

    def _render_image(
        self, prs, layout, spec: SlideSpec, reader: ExcelReader, target: date
    ) -> None:
        from src import image_writer
        if not spec.image_kind:
            return
        png_path = TMP_DIR / f"{spec.image_kind}_{target.isoformat()}.png"
        try:
            image_writer.render(spec.image_kind, reader, target, png_path)
        except image_writer.NoDataError:
            return
        if not png_path.exists():
            return

        slide = prs.slides.add_slide(layout)
        self._set_white_background(slide)
        self._draw_chart_title(slide, spec.title)
        slide.shapes.add_picture(
            str(png_path),
            theme.IMAGE_LEFT, theme.IMAGE_TOP,
            theme.IMAGE_WIDTH, theme.IMAGE_HEIGHT,
        )

    # --- Surplus retrofit helpers ------------------------------------------

    def _inject_surplus_and_total(
        self, rows: list[Row], role: str, reader: ExcelReader, target: date
    ) -> list[Row]:
        """For v1 summary tables: add the surplus row (when sign matches) and
        always append an adjusted-total row."""
        out = list(rows)
        extra = self._surplus_row_for_role(role, reader, target)
        if extra is not None:
            out.append(extra)
        total = sum(r.value for r in out)
        total_label = INCOME_TOTAL_LABEL if role == "income" else EXPENSE_TOTAL_LABEL
        out.append(Row(label=total_label, value=total))
        return out

    @staticmethod
    def _surplus_row_for_role(role: str, reader: ExcelReader,
                              target: date) -> Row | None:
        s = reader.loan_surplus(target)
        if role == "income" and s < 0:
            return Row(label=LOAN_SURPLUS_TITLE_NEGATIVE, value=abs(s))
        if role == "expense" and s > 0:
            return Row(label=LOAN_SURPLUS_TITLE_POSITIVE, value=s)
        return None

    @staticmethod
    def _previous_date(dates: list[date], target: date) -> date | None:
        prev: date | None = None
        for d in sorted(dates):
            if d == target:
                return prev
            prev = d
        return None

    # --- Drawing helpers ---------------------------------------------------

    def _draw_table_title(self, slide, text: str) -> None:
        tb = slide.shapes.add_textbox(
            theme.TABLE_TITLE_LEFT, theme.TABLE_TITLE_TOP,
            theme.TABLE_TITLE_WIDTH, theme.TABLE_TITLE_HEIGHT,
        )
        self._set_textframe_text(
            tb.text_frame, text,
            size=theme.TABLE_TITLE_SIZE, color=theme.RED,
            bold=False, align=PP_ALIGN.CENTER,
        )

    def _draw_chart_title(self, slide, text: str) -> None:
        tb = slide.shapes.add_textbox(
            theme.CHART_TITLE_LEFT, theme.CHART_TITLE_TOP,
            theme.CHART_TITLE_WIDTH, theme.CHART_TITLE_HEIGHT,
        )
        self._set_textframe_text(
            tb.text_frame, text,
            size=theme.CHART_TITLE_SIZE, color=theme.RED,
            bold=False, align=PP_ALIGN.CENTER,
        )

    def _draw_table(self, slide, rows: list[Row]) -> None:
        n_rows = len(rows)
        table_h = theme.TABLE_ROW_HEIGHT * n_rows
        shape = slide.shapes.add_table(
            n_rows, 2,
            theme.TABLE_LEFT, theme.TABLE_TOP,
            theme.TABLE_WIDTH, table_h,
        )
        table = shape.table

        table.first_row = False
        table.first_col = False
        table.horz_banding = False
        table.vert_banding = False

        table.columns[0].width = theme.TABLE_COL0_WIDTH
        table.columns[1].width = theme.TABLE_COL1_WIDTH

        for i, row in enumerate(rows):
            table.rows[i].height = theme.TABLE_ROW_HEIGHT
            bg = theme.TABLE_ROW_A if i % 2 == 0 else theme.TABLE_ROW_B
            self._populate_cell(
                table.cell(i, 0), row.label, bg=bg,
                align=PP_ALIGN.LEFT, bold=False,
                apply_sinhala=True, word_wrap=True,
            )
            self._populate_cell(
                table.cell(i, 1), f"{row.value:,.2f}", bg=bg,
                align=PP_ALIGN.RIGHT, bold=True,
                apply_sinhala=False, word_wrap=False,
            )

    def _draw_n_col_table(
        self,
        slide,
        rows: list[list[str]],
        headers: list[str],
        col_widths: list,
        col_aligns: list,
        col_sinhala: list[bool],
        colorize_columns: dict[int, list] | None = None,
    ) -> None:
        """Generic N-column table for YTD/delta/etc. Includes a header row
        styled in red, then alternating data rows."""
        n_cols = len(headers)
        n_rows = len(rows) + 1   # +1 for header
        # Use a smaller cell-size for these tables; numbers + N columns get tight.
        cell_size = theme.TABLE_CELL_SIZE
        row_h = theme.TABLE_ROW_HEIGHT

        total_w = sum(col_widths)
        shape = slide.shapes.add_table(
            n_rows, n_cols,
            theme.TABLE_LEFT, theme.TABLE_TOP,
            total_w, row_h * n_rows,
        )
        table = shape.table
        table.first_row = False
        table.first_col = False
        table.horz_banding = False
        table.vert_banding = False
        for ci, w in enumerate(col_widths):
            table.columns[ci].width = w

        # Header
        table.rows[0].height = row_h
        for ci, header in enumerate(headers):
            self._populate_cell(
                table.cell(0, ci), header, bg=theme.RED,
                align=PP_ALIGN.CENTER, bold=True,
                apply_sinhala=col_sinhala[ci], word_wrap=True,
                text_color=theme.WHITE, font_size=cell_size,
            )

        # Data
        for ri, row in enumerate(rows, start=1):
            table.rows[ri].height = row_h
            bg = theme.TABLE_ROW_A if (ri - 1) % 2 == 0 else theme.TABLE_ROW_B
            for ci, val in enumerate(row):
                color = theme.BLACK
                if colorize_columns and ci in colorize_columns:
                    custom = colorize_columns[ci][ri - 1]
                    if custom is not None:
                        color = custom
                self._populate_cell(
                    table.cell(ri, ci), val, bg=bg,
                    align=col_aligns[ci],
                    bold=ci > 0,    # numbers bold, label not
                    apply_sinhala=col_sinhala[ci],
                    word_wrap=ci == 0,
                    text_color=color, font_size=cell_size,
                )

    def _populate_cell(self, cell, text: str, bg, align, bold: bool,
                       apply_sinhala: bool, word_wrap: bool,
                       text_color=None, font_size=None) -> None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.12)
        cell.margin_right = Inches(0.12)
        cell.margin_top = Inches(0.05)
        cell.margin_bottom = Inches(0.05)
        tf = cell.text_frame
        tf.clear()
        tf.word_wrap = word_wrap
        para = tf.paragraphs[0]
        para.alignment = align
        run = para.add_run()
        run.text = text
        run.font.size = font_size or theme.TABLE_CELL_SIZE
        run.font.color.rgb = text_color or theme.BLACK
        run.font.bold = bold
        if apply_sinhala:
            apply_sinhala_font(run)
        else:
            run.font.name = theme.LATIN_FONT

    def _set_textframe_text(self, tf, text: str, size, color, bold: bool,
                            align, apply_sinhala: bool = True) -> None:
        tf.clear()
        tf.word_wrap = True
        para = tf.paragraphs[0]
        if align is not None:
            para.alignment = align
        run = para.add_run()
        run.text = text
        if size is not None:
            run.font.size = size
        if color is not None:
            run.font.color.rgb = color
        run.font.bold = bold
        if apply_sinhala:
            apply_sinhala_font(run)
        else:
            run.font.name = theme.LATIN_FONT

    def _fetch_rows(self, query: DataQuery | None, reader: ExcelReader,
                    target: date) -> list[Row]:
        if query is None:
            return []
        col = reader.column_for(query.sheet, target)
        return reader.rows(query.sheet, query.rows, col, query.label_col)

    @staticmethod
    def _set_white_background(slide) -> None:
        """Override the master's dark navy background with explicit white.

        The source deck's master is `#1B2A4A`; only slides 1-3 override it.
        Without this, every new content slide inherits the dark navy.
        """
        cSld = slide.element.find(f"{{{_P_NS}}}cSld")
        existing_bg = cSld.find(f"{{{_P_NS}}}bg")
        if existing_bg is not None:
            cSld.remove(existing_bg)
        bg = etree.SubElement(cSld, f"{{{_P_NS}}}bg")
        bg_pr = etree.SubElement(bg, f"{{{_P_NS}}}bgPr")
        solid = etree.SubElement(bg_pr, f"{{{_A_NS}}}solidFill")
        etree.SubElement(solid, f"{{{_A_NS}}}srgbClr", val="FFFFFF")
        etree.SubElement(bg_pr, f"{{{_A_NS}}}effectLst")
        sp_tree = cSld.find(f"{{{_P_NS}}}spTree")
        if sp_tree is not None and bg.getnext() is not sp_tree:
            cSld.remove(bg)
            cSld.insert(list(cSld).index(sp_tree), bg)

    @staticmethod
    def _cover_placeholders(slide):
        title_ph = None
        subtitle_ph = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                title_ph = ph
            elif ph.placeholder_format.idx == 1:
                subtitle_ph = ph
        return title_ph, subtitle_ph

    @staticmethod
    def _find_layout(prs, name: str):
        for layout in prs.slide_layouts:
            if layout.name == name:
                return layout
        return prs.slide_layouts[-1]

    @staticmethod
    def _remove_slides(prs, sld_ids) -> None:
        sldIdLst = prs.slides._sldIdLst
        for sid in sld_ids:
            rId = sid.get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
            )
            prs.part.drop_rel(rId)
            sldIdLst.remove(sid)


def _distribute_evenly(items, max_per_slide: int) -> list[list]:
    """Split into pages with no slide exceeding `max_per_slide`, distributing
    remainders so page sizes differ by at most 1."""
    n = len(items)
    if n <= max_per_slide:
        return [list(items)]
    n_slides = (n // max_per_slide) + 1
    base = n // n_slides
    extra = n % n_slides
    chunks: list[list] = []
    idx = 0
    for i in range(n_slides):
        size = base + (1 if i < extra else 0)
        chunks.append(list(items[idx:idx + size]))
        idx += size
    return chunks


def _fmt_month(d: date) -> str:
    """Short Sinhala month label for chart axes."""
    months_si = ["", "ජන", "පෙබ", "මාර්", "අප්‍ර", "මැයි", "ජූනි",
                 "ජූලි", "අගෝ", "සැප්", "ඔක්", "නොවැ", "දෙසැ"]
    return f"{months_si[d.month]} {d.year % 100:02d}"
